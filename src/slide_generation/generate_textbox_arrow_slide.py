from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO
from pptx.util import Inches
import random

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

MAX_WIDTH = Inches(10)
MAX_HEIGHT = Inches(7.5)
TOP_PADDING = Inches(2)

MIN_SIZE = Inches(1.5)

def generate_random_grids(nb_grids):
    current_grids = [(0, TOP_PADDING, MAX_WIDTH, MAX_HEIGHT - TOP_PADDING)]
    nb_failures = 0
    current_nb_grids = 1
    while current_nb_grids < nb_grids and nb_failures < 100:
        random_index = random.randint(0, len(current_grids) - 1)
        left, top, width, height = current_grids.pop(random_index)
        if random.randint(0, 1):  # vertical split
            if width < 2 * MIN_SIZE:
                current_grids.append((left, top, width, height))
                nb_failures += 1
                continue
            new_width = random.randint(MIN_SIZE, width - MIN_SIZE)
            other_width = width - new_width
            current_grids.append((left, top, new_width, height))
            current_grids.append((left + new_width, top, other_width, height))
        else:
            if height < 2 * MIN_SIZE:
                current_grids.append((left, top, width, height))
                nb_failures += 1
                continue
            new_height = random.randint(MIN_SIZE, height - MIN_SIZE)
            other_height = height - new_height
            current_grids.append((left, top, width, new_height))
            current_grids.append((left, top + new_height, width, other_height))
        current_nb_grids += 1
    return current_grids



shapes.title.text = 'Title'
grids = generate_random_grids(7)
for left, top, width, height in grids:
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

TEXT_BOX_SIZE = MIN_SIZE // 2
counter = 0
for left, top, width, height in grids:
    vertical_padding = (height - TEXT_BOX_SIZE) // 2
    horizontal_padding = (width - TEXT_BOX_SIZE) // 2
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + horizontal_padding, top + vertical_padding, TEXT_BOX_SIZE, TEXT_BOX_SIZE)
    shape.text = f'Text{counter}'
    counter += 1


prs.save('./test.pptx')